"""
generate_presentations.py
========================
Generate detailed, diagram-rich SQL100 topic presentations for Assmang Pty Ltd.

Run:
    python3 generate_presentations.py

Dependencies:
    pip3 install -r requirements.txt
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path("/Users/kumbulani.tshuma/Documents/my trainning/SQL-Fundamentals-SQL100")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Neutral, reusable palette
C_BG = RGBColor(250, 251, 252)
C_WHITE = RGBColor(255, 255, 255)
C_NAVY = RGBColor(43, 62, 80)
C_BLUE = RGBColor(77, 113, 153)
C_TEAL = RGBColor(88, 145, 164)
C_GREY = RGBColor(94, 109, 126)
C_LIGHT = RGBColor(237, 241, 245)
C_LIGHT_BLUE = RGBColor(222, 233, 245)
C_LIGHT_TEAL = RGBColor(226, 240, 243)
C_LIGHT_GOLD = RGBColor(247, 238, 215)
C_GOLD = RGBColor(181, 137, 63)
C_RED = RGBColor(165, 74, 74)
C_GREEN = RGBColor(48, 112, 76)
C_TEXT = RGBColor(34, 43, 52)
C_MUTED = RGBColor(105, 113, 122)
C_CODE_BG = RGBColor(245, 247, 249)
C_CODE = RGBColor(33, 93, 62)

TITLE_HEIGHT = Inches(0.85)
FOOTER_HEIGHT = Inches(0.32)
CONTENT_TOP = Inches(1.0)
CONTENT_BOTTOM = Inches(7.08)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------
# Drawing helpers
# ---------------------------

def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, *, fill=C_WHITE, line=C_LIGHT, line_width=1.2):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=C_GREY, width=1.4):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_text(slide, text, left, top, width, height, *, size=18, bold=False,
             color=C_TEXT, align=PP_ALIGN.LEFT, font_name="Calibri",
             valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullets(slide, bullets, left, top, width, height, *, size=18, color=C_TEXT, accent=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, str):
            item = {"text": item}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item["text"]
        p.level = item.get("level", 0)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = item.get("color", C_BLUE if accent and p.level == 0 else color)
        if item.get("bullet", True):
            p.bullet = True
        else:
            p.bullet = False
    return box


def add_title_band(slide, title, subtitle=None, day=None):
    set_bg(slide)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_HEIGHT, fill=C_NAVY, line=C_NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.08), Inches(9.5), Inches(0.42), size=26, bold=True, color=C_WHITE, font_name="Calibri Light")
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.45), Inches(9.8), Inches(0.22), size=10.5, color=RGBColor(220, 229, 240))
    if day:
        tag = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(0.12), Inches(2.15), Inches(0.46), fill=C_BLUE, line=C_BLUE)
        add_text(slide, day, Inches(10.82), Inches(0.19), Inches(2.0), Inches(0.18), size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, CONTENT_BOTTOM, SLIDE_W, FOOTER_HEIGHT, fill=C_LIGHT, line=C_LIGHT)
    add_text(slide, "Assmang Pty Ltd | SQL100 | Theory deck", Inches(0.3), Inches(7.1), Inches(5.0), Inches(0.16), size=10, color=C_MUTED)


def title_slide(prs, title, day, slug):
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.95), SLIDE_W, Inches(1.55), fill=C_BLUE, line=C_BLUE)
    add_text(slide, "NobleProg SQL100", Inches(0.5), Inches(0.45), Inches(3.5), Inches(0.3), size=16, bold=True, color=RGBColor(225, 233, 243))
    add_text(slide, day, Inches(0.5), Inches(1.0), Inches(4), Inches(0.4), size=20, color=RGBColor(188, 207, 229))
    add_text(slide, title, Inches(0.5), Inches(1.65), Inches(11.5), Inches(2.4), size=33, bold=True, color=C_WHITE, font_name="Calibri Light")
    add_text(slide, "Detailed theory presentation with diagrams, worked examples, and recap", Inches(0.55), Inches(4.55), Inches(9.5), Inches(0.45), size=15, color=RGBColor(224, 232, 241))
    add_text(slide, slug.replace("-", " ").title(), Inches(0.55), Inches(6.28), Inches(6), Inches(0.26), size=13, color=C_WHITE)
    return slide


def end_slide(prs, title):
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_text(slide, "Questions, Discussion, and Lab Transition", Inches(0.6), Inches(1.4), Inches(11.8), Inches(0.8), size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.2), Inches(2.55), Inches(10.8), Inches(0.7), size=21, color=RGBColor(192, 208, 227), align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.2), Inches(7.3), Inches(1.1), fill=C_BLUE, line=C_BLUE)
    add_text(slide, "Next step: practical lab, guided SQL examples, and discussion questions", Inches(3.2), Inches(4.5), Inches(6.9), Inches(0.3), size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Assmang Pty Ltd | SQL100", Inches(5.0), Inches(6.45), Inches(3.4), Inches(0.2), size=13, color=RGBColor(215, 224, 235), align=PP_ALIGN.CENTER)


def objectives_slide(prs, topic):
    slide = blank_slide(prs)
    add_title_band(slide, "Learning objectives", subtitle=topic["title"], day=topic["day"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(12.3), Inches(5.45), fill=C_WHITE, line=C_LIGHT_BLUE)
    add_text(slide, "By the end of this topic, learners should be able to:", Inches(0.8), Inches(1.55), Inches(6.4), Inches(0.35), size=20, bold=True, color=C_NAVY)
    add_bullets(slide, topic["objectives"], Inches(0.95), Inches(2.05), Inches(7.1), Inches(4.1), size=18, accent=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(2.0), Inches(3.55), Inches(3.4), fill=C_LIGHT_TEAL, line=C_TEAL)
    add_text(slide, "Assmang business lens", Inches(8.85), Inches(2.25), Inches(3.0), Inches(0.28), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_bullets(slide, topic["business_lens"], Inches(8.85), Inches(2.75), Inches(2.85), Inches(2.3), size=15, color=C_TEXT)


def content_slide(prs, title, bullets, *, note=None):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.2), Inches(8.1), Inches(5.55), fill=C_WHITE, line=C_LIGHT_BLUE)
    add_bullets(slide, bullets, Inches(0.8), Inches(1.55), Inches(7.3), Inches(4.8), size=18, accent=True)
    if note:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(1.55), Inches(3.55), Inches(4.9), fill=C_LIGHT_GOLD, line=C_GOLD)
        add_text(slide, "Trainer emphasis", Inches(9.1), Inches(1.85), Inches(3.0), Inches(0.3), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_bullets(slide, note, Inches(9.1), Inches(2.35), Inches(2.95), Inches(3.65), size=15, color=C_TEXT)


def code_slide(prs, title, explanation, code_lines, takeaways=None):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    add_text(slide, explanation, Inches(0.5), Inches(1.08), Inches(12.0), Inches(0.35), size=16, color=C_MUTED)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.55), Inches(7.95), Inches(4.95), fill=C_CODE_BG, line=C_LIGHT_BLUE)
    add_text(slide, "\n".join(code_lines), Inches(0.7), Inches(1.82), Inches(7.45), Inches(4.35), size=15, color=C_CODE, font_name="Courier New")
    if takeaways:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.55), Inches(3.65), Inches(4.95), fill=C_WHITE, line=C_LIGHT_BLUE)
        add_text(slide, "Why this example matters", Inches(9.0), Inches(1.85), Inches(3.0), Inches(0.28), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_bullets(slide, takeaways, Inches(9.0), Inches(2.35), Inches(3.0), Inches(3.8), size=15)


def flow_slide(prs, title, steps, caption):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    step_w = Inches(2.15)
    step_h = Inches(1.55)
    left = Inches(0.6)
    top = Inches(2.2)
    gap = Inches(0.35)
    for idx, step in enumerate(steps):
        x = left + idx * (step_w + gap)
        fill = [C_LIGHT_BLUE, C_LIGHT_TEAL, C_LIGHT_GOLD, C_LIGHT_BLUE, C_LIGHT_TEAL][idx % 5]
        border = [C_BLUE, C_TEAL, C_GOLD, C_BLUE, C_TEAL][idx % 5]
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, step_w, step_h, fill=fill, line=border)
        add_text(slide, step["title"], x + Inches(0.08), top + Inches(0.12), step_w - Inches(0.16), Inches(0.32), size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, step["body"], x + Inches(0.12), top + Inches(0.5), step_w - Inches(0.24), Inches(0.8), size=12.5, color=C_TEXT, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, x + step_w + Inches(0.05), top + Inches(0.42), Inches(0.23), Inches(0.55), fill=C_GREY, line=C_GREY)
    add_text(slide, caption, Inches(0.8), Inches(5.15), Inches(11.7), Inches(0.6), size=17, color=C_MUTED, align=PP_ALIGN.CENTER)


def comparison_slide(prs, title, left_title, left_points, right_title, right_points, footer_text):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.45), Inches(5.8), Inches(4.95), fill=C_WHITE, line=C_BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.45), Inches(5.8), Inches(4.95), fill=C_WHITE, line=C_TEAL)
    add_text(slide, left_title, Inches(1.0), Inches(1.75), Inches(5.0), Inches(0.3), size=21, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, right_title, Inches(7.2), Inches(1.75), Inches(5.1), Inches(0.3), size=21, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_bullets(slide, left_points, Inches(1.0), Inches(2.25), Inches(4.95), Inches(3.55), size=16)
    add_bullets(slide, right_points, Inches(7.25), Inches(2.25), Inches(4.95), Inches(3.55), size=16)
    add_text(slide, footer_text, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.28), size=15, color=C_MUTED, align=PP_ALIGN.CENTER)


def matrix_slide(prs, title, headers, rows, note=None):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    start_x = Inches(0.55)
    start_y = Inches(1.55)
    col_widths = [Inches(2.5), Inches(4.3), Inches(5.1)] if len(headers) == 3 else [Inches(2.1)] * len(headers)
    row_h = Inches(0.58)
    x = start_x
    for idx, header in enumerate(headers):
        w = col_widths[idx]
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, start_y, w, row_h, fill=C_NAVY, line=C_NAVY)
        add_text(slide, header, x + Inches(0.05), start_y + Inches(0.12), w - Inches(0.1), Inches(0.2), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        x += w
    for r_idx, row in enumerate(rows):
        x = start_x
        y = start_y + row_h + r_idx * row_h
        for c_idx, cell in enumerate(row):
            w = col_widths[c_idx]
            fill = C_WHITE if r_idx % 2 == 0 else C_LIGHT
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, row_h, fill=fill, line=C_LIGHT_BLUE, line_width=0.8)
            add_text(slide, str(cell), x + Inches(0.06), y + Inches(0.1), w - Inches(0.12), Inches(0.22), size=12.5, color=C_TEXT)
            x += w
    if note:
        add_text(slide, note, Inches(0.7), Inches(6.55), Inches(11.8), Inches(0.26), size=14, color=C_MUTED, align=PP_ALIGN.CENTER)


def nodes_slide(prs, title, nodes, connections, caption=None):
    slide = blank_slide(prs)
    add_title_band(slide, title)
    box_refs = {}
    for node in nodes:
        x, y, w, h = node["x"], node["y"], node["w"], node["h"]
        box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h,
                        fill=node.get("fill", C_WHITE), line=node.get("line", C_BLUE))
        add_text(slide, node["title"], x + Inches(0.08), y + Inches(0.1), w - Inches(0.16), Inches(0.28), size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_bullets(slide, [{"text": line, "bullet": False, "size": 12.5} for line in node.get("lines", [])],
                    x + Inches(0.15), y + Inches(0.48), w - Inches(0.3), h - Inches(0.58), size=12.5, color=C_TEXT)
        box_refs[node["id"]] = box
    for conn in connections:
        a = next(n for n in nodes if n["id"] == conn["from"])
        b = next(n for n in nodes if n["id"] == conn["to"])
        x1 = a["x"] + a["w"]
        y1 = a["y"] + a["h"] / 2
        x2 = b["x"]
        y2 = b["y"] + b["h"] / 2
        add_line(slide, x1, y1, x2, y2, color=conn.get("color", C_GREY), width=1.5)
        if conn.get("label"):
            lx = (x1 + x2) / 2 - Inches(0.5)
            ly = min(y1, y2) - Inches(0.22)
            add_text(slide, conn["label"], lx, ly, Inches(1.0), Inches(0.2), size=11.5, bold=True, color=conn.get("color", C_GREY), align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, caption, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.24), size=14.5, color=C_MUTED, align=PP_ALIGN.CENTER)


def recap_slide(prs, topic):
    slide = blank_slide(prs)
    add_title_band(slide, "Topic recap", subtitle=topic["title"], day=topic["day"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.45), Inches(5.9), Inches(4.95), fill=C_WHITE, line=C_LIGHT_BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.45), Inches(5.9), Inches(4.95), fill=C_WHITE, line=C_LIGHT_BLUE)
    add_text(slide, "Key messages", Inches(0.85), Inches(1.8), Inches(5.2), Inches(0.3), size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "Common mistakes to avoid", Inches(7.1), Inches(1.8), Inches(5.3), Inches(0.3), size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_bullets(slide, topic["recap"], Inches(0.9), Inches(2.35), Inches(5.0), Inches(3.7), size=16)
    add_bullets(slide, topic["mistakes"], Inches(7.15), Inches(2.35), Inches(5.0), Inches(3.7), size=16)


def render_slide(prs, slide_cfg, topic):
    t = slide_cfg["type"]
    if t == "objectives":
        objectives_slide(prs, topic)
    elif t == "content":
        content_slide(prs, slide_cfg["title"], slide_cfg["bullets"], note=slide_cfg.get("note"))
    elif t == "code":
        code_slide(prs, slide_cfg["title"], slide_cfg["explanation"], slide_cfg["code"], slide_cfg.get("takeaways"))
    elif t == "flow":
        flow_slide(prs, slide_cfg["title"], slide_cfg["steps"], slide_cfg["caption"])
    elif t == "comparison":
        comparison_slide(prs, slide_cfg["title"], slide_cfg["left_title"], slide_cfg["left_points"], slide_cfg["right_title"], slide_cfg["right_points"], slide_cfg["footer_text"])
    elif t == "matrix":
        matrix_slide(prs, slide_cfg["title"], slide_cfg["headers"], slide_cfg["rows"], note=slide_cfg.get("note"))
    elif t == "nodes":
        nodes_slide(prs, slide_cfg["title"], slide_cfg["nodes"], slide_cfg["connections"], slide_cfg.get("caption"))
    elif t == "recap":
        recap_slide(prs, topic)
    else:
        raise ValueError(f"Unknown slide type: {t}")


# ---------------------------
# Topic definitions
# ---------------------------

TOPICS = [
    {
        "day": "Day 1 • Topic 01",
        "title": "Introduction to Databases and SQL",
        "slug": "topic-01-introduction-to-databases-and-sql",
        "output": "Day-01/Topic-01-Introduction-to-Databases-and-SQL/topic-01-introduction-to-databases-and-sql-presentation.pptx",
        "objectives": [
            "Explain what a database is and why organisations use a DBMS instead of spreadsheets alone.",
            "Describe the difference between relational and non-relational storage models.",
            "Recognise tables, rows, columns, keys, and relationships in a relational schema.",
            "Explain what SQL is and identify DQL, DML, DDL, TCL, and DCL at a high level.",
            "Connect the course theory to the Assmang training database and its business entities.",
        ],
        "business_lens": [
            "HR needs employee records.",
            "Operations need mine and equipment data.",
            "Finance needs reliable reports.",
            "Safety needs accurate audit trails.",
        ],
        "recap": [
            "A database is more than storage: it is structure, rules, relationships, and controlled access.",
            "Relational systems are ideal when the business needs consistent reporting and trusted links between tables.",
            "Primary keys identify; foreign keys connect; SQL is the language that works across the entire model.",
        ],
        "mistakes": [
            "Thinking a database is just a bigger spreadsheet with a nicer interface.",
            "Confusing a row with a column, or a primary key with a foreign key.",
            "Using SQL terms without understanding the business meaning behind the data model.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "content", "title": "Why databases exist in business operations", "bullets": [
                "Mining organisations create large volumes of structured data every day: employees, production, maintenance, contractors, and training records.",
                "A database keeps that information consistent, searchable, secure, and shareable across departments.",
                "Compared with unmanaged spreadsheets, a DBMS reduces duplication, supports multiple users, and enforces rules.",
                "For Assmang, this means one trusted source for operational reporting instead of many conflicting versions.",
            ], "note": [
                "Use a real example: employee salary copied into two Excel files can drift apart.",
                "Stress that data quality is a business issue, not only an IT issue.",
            ]},
            {"type": "comparison", "title": "Relational vs non-relational thinking", "left_title": "Relational databases", "left_points": [
                "Data stored in tables with predefined columns.",
                "Relationships are explicit through keys.",
                "Strong fit for payroll, finance, compliance, and reporting.",
                "SQL is the standard query language.",
            ], "right_title": "Non-relational databases", "right_points": [
                "Data stored as documents, key-value pairs, graphs, or wide columns.",
                "Schema can be more flexible and less rigid.",
                "Strong fit for rapidly changing or highly unstructured data.",
                "Often chosen for specialised workloads, not general reporting.",
            ], "footer_text": "For this course, the focus stays on relational SQL because the training data is structured and report-driven."},
            {"type": "nodes", "title": "Core relational building blocks", "nodes": [
                {"id": "table", "title": "Table", "x": Inches(0.8), "y": Inches(1.9), "w": Inches(2.2), "h": Inches(1.8), "fill": C_LIGHT_BLUE, "lines": ["A collection of related records", "Example: employees"]},
                {"id": "row", "title": "Row / Record", "x": Inches(3.6), "y": Inches(1.9), "w": Inches(2.2), "h": Inches(1.8), "fill": C_LIGHT_TEAL, "lines": ["One instance in the table", "Example: one employee"]},
                {"id": "column", "title": "Column / Field", "x": Inches(6.4), "y": Inches(1.9), "w": Inches(2.2), "h": Inches(1.8), "fill": C_LIGHT_GOLD, "lines": ["A single attribute", "Example: hire_date"]},
                {"id": "schema", "title": "Schema", "x": Inches(9.2), "y": Inches(1.9), "w": Inches(2.4), "h": Inches(1.8), "fill": C_LIGHT_BLUE, "lines": ["Blueprint of all tables", "Rules + structure"]},
            ], "connections": [
                {"from": "table", "to": "row", "label": "contains"},
                {"from": "row", "to": "column", "label": "stores values in"},
                {"from": "column", "to": "schema", "label": "defined by"},
            ], "caption": "This vocabulary becomes the foundation for every later SQL statement and diagram."},
            {"type": "nodes", "title": "Primary key and foreign key relationship", "nodes": [
                {"id": "dept", "title": "departments", "x": Inches(1.0), "y": Inches(1.75), "w": Inches(4.0), "h": Inches(2.55), "fill": C_LIGHT_BLUE, "lines": ["department_id  ← primary key", "department_name", "location", "budget_zar"]},
                {"id": "emp", "title": "employees", "x": Inches(7.1), "y": Inches(1.55), "w": Inches(4.6), "h": Inches(3.1), "fill": C_LIGHT_TEAL, "lines": ["employee_id  ← primary key", "first_name / last_name", "department_id  ← foreign key", "salary_zar", "hire_date"]},
            ], "connections": [
                {"from": "dept", "to": "emp", "label": "one department → many employees", "color": C_RED}
            ], "caption": "A foreign key does not replace the primary key; it points back to it so tables stay linked without repeating full department details."},
            {"type": "flow", "title": "Where SQL fits in the data lifecycle", "steps": [
                {"title": "Store", "body": "Data is recorded in tables with rules and types."},
                {"title": "Query", "body": "SQL retrieves only the rows and columns needed."},
                {"title": "Analyse", "body": "Results become reports, dashboards, and decisions."},
                {"title": "Act", "body": "Managers respond using trusted information."},
            ], "caption": "SQL is valuable because it turns stored operational data into useful business answers."},
            {"type": "matrix", "title": "SQL language families at a glance", "headers": ["Family", "Typical commands", "What it means in practice"], "rows": [
                ["DQL", "SELECT", "Ask questions of the data and return results."],
                ["DML", "INSERT, UPDATE, DELETE", "Add, change, or remove table rows."],
                ["DDL", "CREATE, ALTER, DROP", "Define or change database structure."],
                ["TCL", "COMMIT, ROLLBACK", "Control whether changes are saved or undone."],
                ["DCL", "GRANT, REVOKE", "Manage who may use what."],
            ], "note": "Beginners often learn SELECT first, but SQL becomes powerful when they understand where each command family belongs."},
            {"type": "nodes", "title": "Assmang training schema overview", "nodes": [
                {"id": "departments", "title": "departments", "x": Inches(0.9), "y": Inches(1.55), "w": Inches(2.5), "h": Inches(1.55), "fill": C_LIGHT_BLUE, "lines": ["Business units", "Budgets + locations"]},
                {"id": "employees", "title": "employees", "x": Inches(4.15), "y": Inches(1.35), "w": Inches(2.8), "h": Inches(1.95), "fill": C_LIGHT_TEAL, "lines": ["Staff records", "Department + mine links"]},
                {"id": "mines", "title": "mines", "x": Inches(8.05), "y": Inches(1.55), "w": Inches(2.5), "h": Inches(1.55), "fill": C_LIGHT_GOLD, "lines": ["Operations", "Commodity + province"]},
                {"id": "equipment", "title": "equipment", "x": Inches(8.05), "y": Inches(4.05), "w": Inches(2.5), "h": Inches(1.55), "fill": C_LIGHT_BLUE, "lines": ["Trucks, drills, loaders", "Linked to mine"]},
                {"id": "production", "title": "production_monthly", "x": Inches(4.15), "y": Inches(4.05), "w": Inches(2.8), "h": Inches(1.55), "fill": C_LIGHT_TEAL, "lines": ["Tonnes, grade, revenue", "Linked to mine and month"]},
            ], "connections": [
                {"from": "departments", "to": "employees", "label": "department_id"},
                {"from": "employees", "to": "mines", "label": "mine_id"},
                {"from": "mines", "to": "equipment", "label": "mine_id"},
                {"from": "production", "to": "mines", "label": "mine_id"},
            ], "caption": "Later topics use this same schema to show filtering, aggregation, joins, and updates in a realistic business setting."},
            {"type": "code", "title": "First exploratory SQL statements", "explanation": "Early queries should build confidence by proving the learner can connect, inspect, and count data.", "code": [
                "USE assmang_training;",
                "SHOW TABLES;",
                "DESCRIBE employees;",
                "",
                "SELECT * FROM departments;",
                "SELECT * FROM mines;",
                "SELECT COUNT(*) AS total_employees FROM employees;",
                "-- Expected training result: 31",
            ], "takeaways": [
                "SHOW TABLES confirms the schema exists.",
                "DESCRIBE teaches learners to inspect structure before querying.",
                "COUNT(*) gives a fast validation check after data loads.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 1 • Topic 02",
        "title": "SELECT Statements",
        "slug": "topic-02-select-statements",
        "output": "Day-01/Topic-02-SELECT-Statements/topic-02-select-statements-presentation.pptx",
        "objectives": [
            "Read the structure of a basic SELECT query from top to bottom.",
            "Select all columns or only the columns needed for a report.",
            "Use aliases, DISTINCT, and expressions to shape output for business readers.",
            "Recognise the difference between exploratory queries and production-ready queries.",
        ],
        "business_lens": [
            "HR list of staff names.",
            "Finance payroll extracts.",
            "Mine directory reports.",
            "Ad-hoc operational questions.",
        ],
        "recap": [
            "SELECT is about retrieving information, not changing it.",
            "The best SELECT lists only the columns a reader actually needs.",
            "Aliases, DISTINCT, and calculated columns make results easier to consume.",
        ],
        "mistakes": [
            "Using SELECT * in reports where only a few columns are needed.",
            "Forgetting aliases and delivering unclear headings to business users.",
            "Assuming DISTINCT means sorted; it removes duplicates but does not guarantee order.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "flow", "title": "Anatomy of a SELECT statement", "steps": [
                {"title": "SELECT", "body": "Choose columns or expressions to return."},
                {"title": "FROM", "body": "Choose the table that stores the data."},
                {"title": "WHERE", "body": "Optional: filter rows before returning results."},
                {"title": "ORDER BY", "body": "Optional: sort the final result set."},
            ], "caption": "Beginners should see SELECT as a sentence: choose data, say where it lives, optionally filter and sort."},
            {"type": "content", "title": "Why column selection matters", "bullets": [
                "`SELECT *` is useful for quick exploration and validation while learning.",
                "In production reporting, named columns improve readability, performance, and change control.",
                "A narrow query avoids exposing extra fields such as manager IDs or unnecessary salary fields.",
                "Selecting only the needed columns also makes later formatting and joins easier to understand.",
            ], "note": [
                "Tie this to governance: the least information necessary is often the safest option.",
            ]},
            {"type": "code", "title": "Basic SELECT patterns", "explanation": "Start from the simplest useful query patterns before introducing filters or joins.", "code": [
                "SELECT *",
                "FROM employees;",
                "",
                "SELECT first_name, last_name, job_title",
                "FROM employees;",
                "",
                "SELECT department_name, location",
                "FROM departments;",
            ], "takeaways": [
                "The second query is easier to read and more purposeful.",
                "Learners should understand that output columns define what the audience sees.",
            ]},
            {"type": "code", "title": "Aliases make reports readable", "explanation": "Business users respond better to clean column headings than to database field names.", "code": [
                "SELECT",
                "    first_name  AS 'First Name',",
                "    last_name   AS 'Surname',",
                "    job_title   AS 'Position',",
                "    salary_zar  AS 'Monthly Salary (R)'",
                "FROM employees;",
            ], "takeaways": [
                "Aliases do not change table structure; they only change the display label.",
                "Use aliases when technical column names would confuse the audience.",
            ]},
            {"type": "content", "title": "DISTINCT and duplicate thinking", "bullets": [
                "A table can contain many repeated values for the same column, such as job_title or department_id.",
                "`DISTINCT` asks SQL to return each unique value once.",
                "This is useful for exploratory questions like ‘Which mine types exist?’ or ‘Which job titles are in the data?’",
                "When multiple columns are used with DISTINCT, SQL returns unique combinations, not unique values per individual column.",
            ], "note": [
                "Demonstrate how DISTINCT on one column differs from DISTINCT on two columns.",
            ]},
            {"type": "matrix", "title": "Useful SELECT building blocks", "headers": ["Feature", "SQL idea", "Typical beginner use"], "rows": [
                ["All columns", "SELECT *", "Explore a table quickly after loading data."],
                ["Named columns", "SELECT col1, col2", "Build a small, purposeful result set."],
                ["Alias", "AS readable_name", "Present clearer headings in reports."],
                ["DISTINCT", "SELECT DISTINCT col", "List unique values only once."],
                ["Expression", "salary_zar * 12", "Add calculations without changing stored data."],
            ], "note": "This slide works well as a reference screen before practical exercises."},
            {"type": "code", "title": "Expressions in the SELECT list", "explanation": "SQL can calculate values as part of the output without modifying stored data.", "code": [
                "SELECT",
                "    first_name,",
                "    salary_zar AS monthly_salary,",
                "    salary_zar * 12 AS annual_salary,",
                "    ROUND(salary_zar / 22, 2) AS estimated_daily_rate",
                "FROM employees;",
            ], "takeaways": [
                "Expressions are temporary output columns.",
                "Stored salary values do not change unless UPDATE is used later.",
            ]},
            {"type": "comparison", "title": "Exploratory query vs report-ready query", "left_title": "Exploratory", "left_points": [
                "Often uses SELECT *.",
                "Used by analysts while learning or validating data loads.",
                "Fast to write, but output can be cluttered.",
            ], "right_title": "Report-ready", "right_points": [
                "Uses selected columns only.",
                "Adds aliases and business-friendly headings.",
                "Supports cleaner exports to Excel or dashboards.",
            ], "footer_text": "A key beginner skill is knowing when to move from ‘quick check’ SQL to ‘presentation-quality’ SQL."},
            {"type": "code", "title": "Small business-facing example", "explanation": "A simple directory-style query is often more useful than a raw table dump.", "code": [
                "SELECT",
                "    CONCAT(first_name, ' ', last_name) AS employee_name,",
                "    job_title,",
                "    email",
                "FROM employees;",
            ], "takeaways": [
                "This query already feels like a usable contact list.",
                "It introduces expression output without requiring joins yet.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 1 • Topic 03",
        "title": "Filtering with WHERE",
        "slug": "topic-03-filtering-with-where",
        "output": "Day-01/Topic-03-Filtering-with-WHERE/topic-03-filtering-with-where-presentation.pptx",
        "objectives": [
            "Use WHERE to reduce results to only the rows relevant to a question.",
            "Apply comparison operators and logical operators correctly.",
            "Use BETWEEN, IN, LIKE, and NULL checks in realistic scenarios.",
            "Understand why parentheses and NULL handling matter in SQL logic.",
        ],
        "business_lens": [
            "Only safety staff at a site.",
            "Only recent hires.",
            "Only active mines.",
            "Only selected departments.",
        ],
        "recap": [
            "WHERE answers the question: which rows should survive into the result?",
            "AND, OR, and NOT combine conditions but require careful grouping.",
            "NULL needs `IS NULL` / `IS NOT NULL`, not `= NULL`.",
        ],
        "mistakes": [
            "Using `= NULL` instead of `IS NULL`.",
            "Mixing OR and AND without parentheses.",
            "Using LIKE when an exact match with `=` would be clearer and safer.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "flow", "title": "How WHERE changes a result set", "steps": [
                {"title": "Full table", "body": "Start with all rows in the source table."},
                {"title": "Condition", "body": "Evaluate each row against the rule."},
                {"title": "Keep / remove", "body": "Rows that pass stay; others disappear."},
                {"title": "Return", "body": "SELECT shows only the surviving rows."},
            ], "caption": "Filtering happens before results are presented, which is why WHERE is so important for business accuracy."},
            {"type": "matrix", "title": "Comparison operators", "headers": ["Operator", "Meaning", "Assmang example"], "rows": [
                ["=", "Exactly equal", "department_id = 2"],
                ["<>", "Not equal", "mine_type <> 'Chrome'"],
                ["> / <", "Greater than / less than", "salary_zar > 100000"],
                [">= / <=", "Inclusive comparisons", "hire_date >= '2020-01-01'"],
                ["BETWEEN", "Range including both ends", "salary_zar BETWEEN 50000 AND 80000"],
            ], "note": "Use a quick verbal reminder that date comparisons still work because SQL stores dates as proper date values, not loose text."},
            {"type": "code", "title": "Simple WHERE examples", "explanation": "Learners should see a direct link between a business question and a single filter condition.", "code": [
                "SELECT first_name, last_name, salary_zar",
                "FROM employees",
                "WHERE salary_zar > 100000;",
                "",
                "SELECT mine_name, mine_type",
                "FROM mines",
                "WHERE mine_type = 'Iron Ore';",
            ], "takeaways": [
                "The filter belongs after FROM and before ORDER BY.",
                "Each condition should read like a clear rule.",
            ]},
            {"type": "comparison", "title": "AND vs OR thinking", "left_title": "AND", "left_points": [
                "Every condition must be true.",
                "Narrows the result set.",
                "Useful for ‘employees in engineering AND earning above R80,000’."
            ], "right_title": "OR", "right_points": [
                "At least one condition must be true.",
                "Widens the result set.",
                "Useful for ‘Chrome OR Manganese mines’."
            ], "footer_text": "When AND and OR appear together, parentheses protect the business meaning of the query."},
            {"type": "code", "title": "Combining logical conditions safely", "explanation": "Parentheses make operator precedence explicit and reduce accidental misinterpretation.", "code": [
                "SELECT first_name, last_name, department_id, salary_zar",
                "FROM employees",
                "WHERE (department_id = 2 OR department_id = 3)",
                "  AND salary_zar > 60000;",
            ], "takeaways": [
                "Without parentheses, SQL applies AND before OR.",
                "Always write the logic the way you would explain it aloud.",
            ]},
            {"type": "content", "title": "BETWEEN, IN, and LIKE are convenience tools", "bullets": [
                "`BETWEEN` is best when a value must fall inside a range, such as salary bands or date windows.",
                "`IN` reads well when a field may match one of several listed values.",
                "`LIKE` supports pattern matching with `%` for any number of characters and `_` for exactly one character.",
                "The simpler the condition, the easier it is to explain, debug, and teach.",
            ], "note": [
                "Show that `IN (1,4,6)` is easier to read than three OR conditions.",
            ]},
            {"type": "code", "title": "Pattern and list matching examples", "explanation": "These patterns appear regularly in admin, HR, and operations filters.", "code": [
                "SELECT first_name, last_name, job_title",
                "FROM employees",
                "WHERE job_title IN ('Driller', 'Blaster', 'Truck Operator');",
                "",
                "SELECT mine_name",
                "FROM mines",
                "WHERE mine_name LIKE 'B%';",
            ], "takeaways": [
                "IN improves readability for lists.",
                "LIKE is valuable, but it should not replace exact matching without reason.",
            ]},
            {"type": "content", "title": "NULL changes the logic model", "bullets": [
                "NULL means ‘unknown’ or ‘not currently stored’, not zero and not an empty string.",
                "Because NULL is not a normal value, `= NULL` and `<> NULL` do not behave as beginners expect.",
                "Use `IS NULL` to find missing data and `IS NOT NULL` to find stored data.",
                "In the Assmang dataset, a NULL `mine_id` identifies head-office staff rather than site-based staff.",
            ], "note": [
                "This is often the first major logic surprise for SQL beginners.",
            ]},
            {"type": "code", "title": "Correct NULL handling", "explanation": "Null filters are essential in real reporting because optional relationships are common.", "code": [
                "SELECT first_name, last_name, job_title",
                "FROM employees",
                "WHERE mine_id IS NULL;",
                "",
                "SELECT first_name, last_name, mine_id",
                "FROM employees",
                "WHERE mine_id IS NOT NULL;",
            ], "takeaways": [
                "This pattern separates head-office staff from mine-based staff.",
                "It also prepares learners for LEFT JOIN logic later in the course.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 1 • Topic 04",
        "title": "Sorting and Limiting Results",
        "slug": "topic-04-sorting-and-limiting-results",
        "output": "Day-01/Topic-04-Sorting-and-Limiting-Results/topic-04-sorting-and-limiting-results-presentation.pptx",
        "objectives": [
            "Sort SQL results using ORDER BY with ascending and descending logic.",
            "Apply multi-column sorting to produce stable, readable outputs.",
            "Use LIMIT and OFFSET for top-N lists and basic pagination.",
            "Explain how ordering fits into the wider query execution flow.",
        ],
        "business_lens": [
            "Top earners report.",
            "Latest hires list.",
            "Paginated staff directory.",
            "Budget ranking outputs.",
        ],
        "recap": [
            "ORDER BY changes presentation order, not stored data.",
            "Sort order should be intentional, not accidental.",
            "LIMIT is most meaningful when combined with an explicit sort.",
        ],
        "mistakes": [
            "Using LIMIT without ORDER BY and assuming the ‘top’ rows are meaningful.",
            "Forgetting that secondary sorts stabilise rows with equal primary values.",
            "Confusing OFFSET counts with page numbers.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "content", "title": "Why sorting matters to business users", "bullets": [
                "A result set without ORDER BY may be technically correct but hard to interpret.",
                "Sorted outputs support ranking, comparison, and quick decision-making.",
                "Managers often think in ordered lists: highest revenue, newest hire, lowest stock, oldest incident.",
                "The database stores data one way; ORDER BY lets us present it a different way for the reader.",
            ], "note": [
                "Reinforce that sorted output is a presentation feature, not a change to stored table order.",
            ]},
            {"type": "comparison", "title": "Ascending and descending order", "left_title": "ASC", "left_points": [
                "Default sort direction in SQL.",
                "Useful for A→Z, oldest→newest, lowest→highest.",
                "Great for employee directories or chronological logs.",
            ], "right_title": "DESC", "right_points": [
                "Explicit reverse ordering.",
                "Useful for highest→lowest, newest→oldest, top performers first.",
                "Great for dashboards and leadership summaries.",
            ], "footer_text": "Use ASC and DESC deliberately so the reader instantly sees the intended ranking."},
            {"type": "code", "title": "Basic ORDER BY examples", "explanation": "These queries show how the same data becomes more useful once ordered.", "code": [
                "SELECT first_name, last_name, salary_zar",
                "FROM employees",
                "ORDER BY salary_zar DESC;",
                "",
                "SELECT mine_name, established_year",
                "FROM mines",
                "ORDER BY established_year ASC;",
            ], "takeaways": [
                "Descending salary produces an immediate ranking view.",
                "Ascending year helps readers understand historical sequence.",
            ]},
            {"type": "flow", "title": "Multi-column sort logic", "steps": [
                {"title": "Primary sort", "body": "Group rows by the most important ordering field."},
                {"title": "Tie handling", "body": "Rows with equal primary values remain ambiguous."},
                {"title": "Secondary sort", "body": "A second column resolves the tie consistently."},
                {"title": "Stable output", "body": "The final list is easier to read and reuse."},
            ], "caption": "Multi-column ORDER BY is the SQL equivalent of saying ‘first sort by department, then within each department sort by salary’."},
            {"type": "code", "title": "Multi-column ordering", "explanation": "This is one of the most useful reporting patterns for grouped business lists.", "code": [
                "SELECT department_id, last_name, salary_zar",
                "FROM employees",
                "ORDER BY department_id ASC, salary_zar DESC;",
            ], "takeaways": [
                "Readers can review one department at a time.",
                "Within each department, the highest-paid roles appear first.",
            ]},
            {"type": "flow", "title": "LIMIT and OFFSET for pagination", "steps": [
                {"title": "Sort", "body": "Choose a stable order, such as employee_id or surname."},
                {"title": "Skip", "body": "OFFSET ignores the earlier rows."},
                {"title": "Return", "body": "LIMIT returns only the next set of rows."},
                {"title": "Repeat", "body": "The next page changes only the OFFSET."},
            ], "caption": "Pagination works best when the sort is predictable; otherwise rows may appear to ‘move’ between pages."},
            {"type": "code", "title": "Top-N and page examples", "explanation": "Two common patterns are the ‘top list’ and the paginated list.", "code": [
                "SELECT first_name, last_name, salary_zar",
                "FROM employees",
                "ORDER BY salary_zar DESC",
                "LIMIT 5;",
                "",
                "SELECT employee_id, first_name, last_name",
                "FROM employees",
                "ORDER BY employee_id",
                "LIMIT 5 OFFSET 10;",
            ], "takeaways": [
                "The first query answers ‘Who are the top 5 earners?’",
                "The second query answers ‘Show me page 3 when each page contains 5 rows’."
            ]},
            {"type": "content", "title": "ORDER BY in the execution model", "bullets": [
                "SQL does not sort first; it first determines the source rows and any filters.",
                "Only after the result set is assembled does ORDER BY change how that set is displayed.",
                "This is why aliases created in SELECT can often be used in ORDER BY even though they are not available in WHERE.",
                "Understanding sequence prevents many beginner mistakes later in GROUP BY and HAVING.",
            ], "note": [
                "Link this to the next day: learners need this ordering intuition before aggregation.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 1 • Topic 05",
        "title": "Single-Row Functions",
        "slug": "topic-05-single-row-functions",
        "output": "Day-01/Topic-05-Single-Row-Functions/topic-05-single-row-functions-presentation.pptx",
        "objectives": [
            "Recognise the main categories of single-row functions: string, numeric, date, conversion, and conditional.",
            "Use functions to transform values without changing stored data.",
            "Apply CASE, COALESCE, and formatting functions in business-facing queries.",
            "Choose functions that clarify results instead of making queries harder to read.",
        ],
        "business_lens": [
            "Readable names and emails.",
            "Salary calculations.",
            "Years of service.",
            "Status labels for reports.",
        ],
        "recap": [
            "Single-row functions work on one row at a time and return one value per row.",
            "Functions improve readability, classification, and formatting for business outputs.",
            "CASE and COALESCE are especially powerful when data needs interpretation, not just retrieval.",
        ],
        "mistakes": [
            "Using functions everywhere and making simple SQL hard to read.",
            "Confusing formatting logic with stored business logic.",
            "Forgetting that function output is temporary unless data is updated later.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "comparison", "title": "Single-row vs aggregate functions", "left_title": "Single-row", "left_points": [
                "Processes one row at a time.",
                "Returns one output value per input row.",
                "Examples: UPPER, ROUND, YEAR, COALESCE, CASE.",
            ], "right_title": "Aggregate", "right_points": [
                "Processes multiple rows together.",
                "Returns one summary value for a set or group.",
                "Examples: COUNT, SUM, AVG, MIN, MAX.",
            ], "footer_text": "This topic stays at the row level; grouping arrives in Day 2 Topic 06."},
            {"type": "content", "title": "Function categories learners should recognise", "bullets": [
                "String functions shape text values into cleaner outputs.",
                "Numeric functions handle rounding, absolute values, and simple mathematical cleanup.",
                "Date functions interpret dates and calculate service periods or reporting windows.",
                "Conditional and null-handling functions help SQL express business meaning, not just raw values.",
            ], "note": [
                "Frame functions as tools for presentation and interpretation.",
            ]},
            {"type": "code", "title": "String functions", "explanation": "String functions are often the fastest way to make output more readable to humans.", "code": [
                "SELECT",
                "    UPPER(last_name) AS surname_caps,",
                "    LOWER(email) AS email_lower,",
                "    CONCAT(first_name, ' ', last_name) AS full_name,",
                "    LENGTH(job_title) AS title_length,",
                "    REPLACE(email, '@assmang.co.za', '') AS username",
                "FROM employees;",
            ], "takeaways": [
                "These functions affect only the result set, not the stored table values.",
                "String cleanup is often necessary before exporting reports.",
            ]},
            {"type": "code", "title": "Numeric functions", "explanation": "Numeric functions are useful whenever raw values need adjustment or presentation.", "code": [
                "SELECT",
                "    salary_zar,",
                "    ROUND(salary_zar / 22, 2) AS estimated_daily_rate,",
                "    FLOOR(salary_zar / 1000) AS salary_thousands_floor,",
                "    ABS(salary_zar - 50000) AS distance_from_band",
                "FROM employees;",
            ], "takeaways": [
                "ROUND is common in payroll-style reports.",
                "ABS helps compare distance from a target value.",
            ]},
            {"type": "code", "title": "Date functions", "explanation": "Date logic often matters more to the business than the raw stored date itself.", "code": [
                "SELECT",
                "    hire_date,",
                "    YEAR(hire_date) AS hire_year,",
                "    DATE_FORMAT(hire_date, '%d %M %Y') AS formatted_hire_date,",
                "    TIMESTAMPDIFF(YEAR, hire_date, CURRENT_DATE()) AS years_service",
                "FROM employees;",
            ], "takeaways": [
                "Formatting makes dates friendlier to readers.",
                "TIMESTAMPDIFF turns a date into business context such as years of service.",
            ]},
            {"type": "flow", "title": "CASE translates values into business language", "steps": [
                {"title": "Read row", "body": "Evaluate the current salary or status value."},
                {"title": "Test rule 1", "body": "If the first condition matches, return its label."},
                {"title": "Test next", "body": "If not, continue through later conditions."},
                {"title": "Return label", "body": "Convert raw data into a business-friendly category."},
            ], "caption": "CASE is one of the most powerful beginner tools because it converts data values into reporting language."},
            {"type": "code", "title": "CASE and COALESCE", "explanation": "These functions help SQL explain the data instead of only listing the data.", "code": [
                "SELECT",
                "    first_name,",
                "    salary_zar,",
                "    CASE",
                "        WHEN salary_zar >= 100000 THEN 'Executive'",
                "        WHEN salary_zar >= 75000 THEN 'Senior'",
                "        WHEN salary_zar >= 50000 THEN 'Mid-Level'",
                "        ELSE 'Junior'",
                "    END AS salary_band,",
                "    COALESCE(CAST(mine_id AS CHAR), 'Head Office') AS site",
                "FROM employees;",
            ], "takeaways": [
                "CASE categorises raw values into business tiers.",
                "COALESCE provides a fallback when a value is NULL.",
            ]},
            {"type": "content", "title": "Choosing functions responsibly", "bullets": [
                "A function should improve clarity, not create complexity for its own sake.",
                "If several nested functions are hard to explain aloud, the query may need simplification.",
                "Well-chosen functions make exports and dashboards easier to understand without changing the underlying data model.",
                "In beginner SQL, readability is as important as correctness.",
            ], "note": [
                "Encourage learners to write SQL they can explain line by line.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 2 • Topic 06",
        "title": "Aggregate Functions and GROUP BY",
        "slug": "topic-06-aggregate-functions-and-group-by",
        "output": "Day-02/Topic-06-Aggregate-Functions-and-GROUP-BY/topic-06-aggregate-functions-and-group-by-presentation.pptx",
        "objectives": [
            "Use COUNT, SUM, AVG, MIN, and MAX to summarise data sets.",
            "Explain the difference between row-level queries and grouped summary queries.",
            "Use GROUP BY to create business summaries per department, mine, or status.",
            "Use HAVING to filter groups after aggregation.",
        ],
        "business_lens": [
            "Headcount by department.",
            "Payroll by team.",
            "Production by mine.",
            "Equipment value summaries.",
        ],
        "recap": [
            "Aggregate functions summarise sets; GROUP BY decides how those sets are formed.",
            "WHERE filters rows before grouping; HAVING filters the grouped results.",
            "Good aggregation turns raw operational data into managerial insight.",
        ],
        "mistakes": [
            "Mixing non-grouped columns into a grouped SELECT list.",
            "Trying to use aggregate logic in WHERE instead of HAVING.",
            "Confusing COUNT(*) with COUNT(column) when NULL values exist.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "comparison", "title": "Row-level thinking vs summary-level thinking", "left_title": "Row-level query", "left_points": [
                "Returns one row per stored record.",
                "Used when the reader needs detail.",
                "Example: each employee and their salary.",
            ], "right_title": "Summary-level query", "right_points": [
                "Returns one row per group or one row overall.",
                "Used when the reader needs totals or averages.",
                "Example: payroll per department.",
            ], "footer_text": "Aggregation matters because managers usually want patterns and totals, not raw row dumps."},
            {"type": "matrix", "title": "Aggregate functions at a glance", "headers": ["Function", "Question it answers", "Example use"], "rows": [
                ["COUNT(*)", "How many rows are there?", "How many employees exist?"],
                ["SUM(col)", "What is the total?", "What is the monthly payroll total?"],
                ["AVG(col)", "What is the average?", "What is the average salary?"],
                ["MIN(col)", "What is the smallest value?", "What is the earliest hire date?"],
                ["MAX(col)", "What is the largest value?", "What is the highest revenue month?"],
            ], "note": "The teaching goal is not memorisation alone, but recognising which business question maps to which function."},
            {"type": "code", "title": "Whole-table aggregation", "explanation": "These functions are easiest to learn when first applied to the full table.", "code": [
                "SELECT",
                "    COUNT(*) AS total_employees,",
                "    ROUND(AVG(salary_zar), 2) AS avg_salary,",
                "    MIN(salary_zar) AS lowest_salary,",
                "    MAX(salary_zar) AS highest_salary,",
                "    SUM(salary_zar) AS total_monthly_payroll",
                "FROM employees;",
            ], "takeaways": [
                "One input table can become one summary row.",
                "This is the foundation for grouped summaries next.",
            ]},
            {"type": "flow", "title": "How GROUP BY forms summaries", "steps": [
                {"title": "Read rows", "body": "Start with the filtered table rows."},
                {"title": "Partition", "body": "Split rows into groups using the chosen column."},
                {"title": "Aggregate", "body": "Run COUNT, SUM, AVG, etc. inside each group."},
                {"title": "Return", "body": "Produce one output row per group."},
            ], "caption": "GROUP BY does not merely sort data; it changes the grain of the result from detail to summary."},
            {"type": "code", "title": "Grouped payroll summary", "explanation": "This is the kind of query managers often expect after learning basic SQL.", "code": [
                "SELECT",
                "    department_id,",
                "    COUNT(*) AS headcount,",
                "    ROUND(AVG(salary_zar), 2) AS avg_salary,",
                "    SUM(salary_zar) AS monthly_payroll",
                "FROM employees",
                "GROUP BY department_id",
                "ORDER BY monthly_payroll DESC;",
            ], "takeaways": [
                "Every selected non-aggregate column must appear in GROUP BY.",
                "The result grain is now ‘one row per department’."
            ]},
            {"type": "comparison", "title": "WHERE vs HAVING", "left_title": "WHERE", "left_points": [
                "Filters individual rows.",
                "Runs before GROUP BY.",
                "Used for conditions like is_active = 1.",
            ], "right_title": "HAVING", "right_points": [
                "Filters grouped results.",
                "Runs after aggregation.",
                "Used for conditions like COUNT(*) >= 3.",
            ], "footer_text": "A simple memory rule: WHERE chooses rows, HAVING chooses groups."},
            {"type": "code", "title": "HAVING in action", "explanation": "HAVING is how grouped reports become targeted management reports.", "code": [
                "SELECT department_id, COUNT(*) AS headcount",
                "FROM employees",
                "GROUP BY department_id",
                "HAVING COUNT(*) >= 3",
                "ORDER BY headcount DESC;",
            ], "takeaways": [
                "The query first builds department groups, then removes the small groups.",
                "This is not possible with WHERE because the counts do not yet exist there.",
            ]},
            {"type": "flow", "title": "Execution order for grouped queries", "steps": [
                {"title": "FROM", "body": "Identify source table(s)."},
                {"title": "WHERE", "body": "Filter row-level conditions."},
                {"title": "GROUP BY", "body": "Build groups."},
                {"title": "HAVING", "body": "Keep or discard groups."},
                {"title": "SELECT/ORDER", "body": "Return and sort the grouped result."},
            ], "caption": "This sequence is critical for explaining why some aliases or aggregates work in one clause but not another."},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 2 • Topic 07",
        "title": "Joining Tables",
        "slug": "topic-07-joining-tables",
        "output": "Day-02/Topic-07-Joining-Tables/topic-07-joining-tables-presentation.pptx",
        "objectives": [
            "Explain why related data is stored across multiple tables instead of being duplicated.",
            "Use INNER, LEFT, and RIGHT joins appropriately.",
            "Read join conditions clearly and connect them to foreign keys.",
            "Use aliases and multi-table joins to create richer business reports.",
        ],
        "business_lens": [
            "Employee + department report.",
            "Employee + mine assignment.",
            "Equipment by mine.",
            "Departmental payroll directory.",
        ],
        "recap": [
            "JOINs reassemble related information that was intentionally split across tables.",
            "The join type changes which unmatched rows are kept or discarded.",
            "Good joins are built on clear keys and readable aliases.",
        ],
        "mistakes": [
            "Forgetting the join condition and producing a Cartesian product.",
            "Using INNER JOIN when unmatched rows should still be shown.",
            "Selecting columns without table aliases when names overlap across tables.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "content", "title": "Why relational data is split across tables", "bullets": [
                "A well-designed schema avoids repeating long descriptive values in every row.",
                "Instead of storing department name repeatedly inside employees, the schema stores department_id once per employee and the department details in a separate table.",
                "This reduces duplication, improves maintainability, and prevents contradictory copies of the same business fact.",
                "JOINs are how SQL temporarily rebuilds the richer view needed by a report or dashboard.",
            ], "note": [
                "Explain normalisation in simple language: store the fact once, then link to it.",
            ]},
            {"type": "nodes", "title": "Relationship path used in many reports", "nodes": [
                {"id": "d", "title": "departments", "x": Inches(0.9), "y": Inches(2.1), "w": Inches(2.8), "h": Inches(1.7), "fill": C_LIGHT_BLUE, "lines": ["department_id", "department_name", "location"]},
                {"id": "e", "title": "employees", "x": Inches(5.2), "y": Inches(1.75), "w": Inches(2.8), "h": Inches(2.4), "fill": C_LIGHT_TEAL, "lines": ["employee_id", "department_id", "mine_id", "job_title"]},
                {"id": "m", "title": "mines", "x": Inches(9.4), "y": Inches(2.1), "w": Inches(2.6), "h": Inches(1.7), "fill": C_LIGHT_GOLD, "lines": ["mine_id", "mine_name", "mine_type"]},
            ], "connections": [
                {"from": "d", "to": "e", "label": "department_id"},
                {"from": "e", "to": "m", "label": "mine_id"},
            ], "caption": "A single employee report may need information from three tables even though no table stores every field together."},
            {"type": "comparison", "title": "Join type behaviour", "left_title": "INNER JOIN", "left_points": [
                "Keep rows only when both tables match.",
                "Best when only complete, matched records matter.",
                "Head-office staff disappear from employee↔mine joins because mine_id is NULL.",
            ], "right_title": "LEFT JOIN", "right_points": [
                "Keep all rows from the left table.",
                "Unmatched right-side values become NULL.",
                "Best when the left table is the primary business list.",
            ], "footer_text": "RIGHT JOIN is conceptually similar to LEFT JOIN, but with table sides reversed."},
            {"type": "code", "title": "Employee plus department using INNER JOIN", "explanation": "This is the first classic reporting join in most SQL courses.", "code": [
                "SELECT",
                "    e.first_name,",
                "    e.last_name,",
                "    e.job_title,",
                "    d.department_name",
                "FROM employees e",
                "INNER JOIN departments d",
                "    ON e.department_id = d.department_id;",
            ], "takeaways": [
                "Aliases keep the query readable.",
                "The ON clause states exactly how the tables relate.",
            ]},
            {"type": "code", "title": "Employee plus mine using LEFT JOIN", "explanation": "LEFT JOIN is essential when some employees do not belong to a mine site.", "code": [
                "SELECT",
                "    CONCAT(e.first_name, ' ', e.last_name) AS employee,",
                "    e.job_title,",
                "    COALESCE(m.mine_name, 'Head Office') AS site",
                "FROM employees e",
                "LEFT JOIN mines m",
                "    ON e.mine_id = m.mine_id;",
            ], "takeaways": [
                "Every employee remains in the result.",
                "COALESCE converts the NULL site into a readable business label.",
            ]},
            {"type": "flow", "title": "How to think through a join question", "steps": [
                {"title": "Choose base list", "body": "Which table contains the main rows you care about?"},
                {"title": "Identify link", "body": "Which key connects to the second table?"},
                {"title": "Pick join type", "body": "Do unmatched base rows need to remain visible?"},
                {"title": "Name outputs", "body": "Choose columns that answer the business question."},
            ], "caption": "Teaching join logic as a decision process helps learners avoid random trial-and-error queries."},
            {"type": "code", "title": "Three-table join pattern", "explanation": "Multi-table joins create the rich outputs usually expected in real business reporting.", "code": [
                "SELECT",
                "    CONCAT(e.first_name, ' ', e.last_name) AS employee,",
                "    d.department_name,",
                "    COALESCE(m.mine_name, 'Head Office') AS site",
                "FROM employees e",
                "INNER JOIN departments d ON e.department_id = d.department_id",
                "LEFT JOIN mines m ON e.mine_id = m.mine_id;",
            ], "takeaways": [
                "The result now answers three questions at once: who, which department, and which site.",
            ]},
            {"type": "content", "title": "Readability practices for joins", "bullets": [
                "Use short aliases such as e, d, and m to reduce visual clutter.",
                "Always qualify shared column names, especially ids and names, to avoid ambiguity.",
                "Format each JOIN and ON clause on separate lines so learners can visually trace the data path.",
                "Explain joins from the business point of view before diving into syntax details.",
            ], "note": [
                "For beginners, formatting is part of understanding.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 2 • Topic 08",
        "title": "Subqueries",
        "slug": "topic-08-subqueries",
        "output": "Day-02/Topic-08-Subqueries/topic-08-subqueries-presentation.pptx",
        "objectives": [
            "Explain what a subquery is and why it can be useful.",
            "Use subqueries in WHERE, IN, EXISTS, FROM, and SELECT positions.",
            "Recognise when a subquery is clearer than a join and when it is not.",
            "Understand the basic idea of correlated subqueries.",
        ],
        "business_lens": [
            "Above-average earners.",
            "Employees in selected departments.",
            "Mines with or without production.",
            "Derived summary filters.",
        ],
        "recap": [
            "A subquery is a query inside another query.",
            "Subqueries are strong when the outer query depends on a lookup, existence test, or derived summary.",
            "Choosing between a join and a subquery is often about readability and intent.",
        ],
        "mistakes": [
            "Using `=` when the subquery may return multiple rows.",
            "Building a subquery when a simple join would be easier to explain.",
            "Treating correlated subqueries as magic instead of row-by-row logic.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "flow", "title": "Mental model for subqueries", "steps": [
                {"title": "Inner query", "body": "Find a value, set, or existence result."},
                {"title": "Return result", "body": "Pass that result outward."},
                {"title": "Outer query", "body": "Use the returned result to filter or display data."},
                {"title": "Final answer", "body": "Return the rows the business actually asked for."},
            ], "caption": "Subqueries are easiest to understand when treated as helper questions inside a bigger business question."},
            {"type": "code", "title": "Single-value subquery in WHERE", "explanation": "A scalar subquery returns one value used by the outer WHERE clause.", "code": [
                "SELECT first_name, last_name, salary_zar",
                "FROM employees",
                "WHERE salary_zar > (",
                "    SELECT AVG(salary_zar)",
                "    FROM employees",
                ");",
            ], "takeaways": [
                "The inner query answers: what is the average salary?",
                "The outer query then asks: who is above that average?",
            ]},
            {"type": "comparison", "title": "`=` vs `IN` with subqueries", "left_title": "Use =", "left_points": [
                "The subquery returns exactly one value.",
                "Example: compare to a single average salary.",
                "Good for scalar comparisons.",
            ], "right_title": "Use IN", "right_points": [
                "The subquery may return many values.",
                "Example: all department_ids located in Johannesburg.",
                "Good when matching against a result list.",
            ], "footer_text": "This distinction prevents one of the most common beginner runtime errors in subquery work."},
            {"type": "code", "title": "Multi-row subquery with IN", "explanation": "IN works well when the inner query produces a list of acceptable values.", "code": [
                "SELECT first_name, last_name, department_id",
                "FROM employees",
                "WHERE department_id IN (",
                "    SELECT department_id",
                "    FROM departments",
                "    WHERE location LIKE '%Johannesburg%'",
                ");",
            ], "takeaways": [
                "The departments table decides the valid ids.",
                "The employees table uses that returned set as a filter.",
            ]},
            {"type": "content", "title": "EXISTS and NOT EXISTS", "bullets": [
                "`EXISTS` checks whether the inner query returns at least one row.",
                "It is often the clearest way to ask ‘Does a related record exist?’",
                "`NOT EXISTS` is the opposite: show rows for which the related record does not exist.",
                "These patterns are especially useful in audits, validation checks, and exception reports.",
            ], "note": [
                "This is a strong bridge between joins and data quality thinking.",
            ]},
            {"type": "code", "title": "EXISTS example", "explanation": "This pattern asks about existence without needing to list fields from the inner table.", "code": [
                "SELECT mine_name, mine_type",
                "FROM mines m",
                "WHERE EXISTS (",
                "    SELECT 1",
                "    FROM production_monthly p",
                "    WHERE p.mine_id = m.mine_id",
                ");",
            ], "takeaways": [
                "`SELECT 1` is enough because the existence of rows matters, not their detailed values.",
            ]},
            {"type": "code", "title": "Subquery in FROM (derived table)", "explanation": "Sometimes the outer query needs to use a summary as though it were a temporary table.", "code": [
                "SELECT summary.department_id, summary.avg_salary",
                "FROM (",
                "    SELECT department_id, AVG(salary_zar) AS avg_salary",
                "    FROM employees",
                "    GROUP BY department_id",
                ") AS summary",
                "WHERE summary.avg_salary > 60000;",
            ], "takeaways": [
                "The inner query builds a temporary summary table.",
                "The outer query filters that derived table.",
            ]},
            {"type": "flow", "title": "Correlated subquery idea", "steps": [
                {"title": "Outer row", "body": "Pick one row from the outer query."},
                {"title": "Inner check", "body": "Run a subquery that refers back to that outer row."},
                {"title": "Compare", "body": "Decide whether the row passes the rule."},
                {"title": "Repeat", "body": "Do it again for every outer row."},
            ], "caption": "A correlated subquery is evaluated per outer row, which is why it feels different from a one-time subquery."},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 2 • Topic 09",
        "title": "DML — INSERT, UPDATE, DELETE",
        "slug": "topic-09-dml-insert-update-delete",
        "output": "Day-02/Topic-09-DML-INSERT-UPDATE-DELETE/topic-09-dml-insert-update-delete-presentation.pptx",
        "objectives": [
            "Use INSERT to add single and multiple rows safely.",
            "Use UPDATE and DELETE with careful WHERE filtering.",
            "Understand transaction safety with START TRANSACTION, COMMIT, and ROLLBACK.",
            "Explain the difference between TRUNCATE and DELETE at a beginner level.",
        ],
        "business_lens": [
            "New hires.",
            "Salary updates.",
            "Removing bad training rows.",
            "Undoing mistakes safely.",
        ],
        "recap": [
            "DML changes data, so safety and verification matter more than speed.",
            "The safest habit is: SELECT first, then UPDATE or DELETE using the same WHERE logic.",
            "Transactions create a controlled window in which mistakes can still be reversed.",
        ],
        "mistakes": [
            "Running UPDATE or DELETE without WHERE.",
            "Changing data before verifying which rows will be affected.",
            "Using TRUNCATE as though it were just a faster DELETE.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "flow", "title": "Safe DML operating pattern", "steps": [
                {"title": "Inspect", "body": "Run SELECT to confirm target rows."},
                {"title": "Begin", "body": "Open a transaction for risky changes."},
                {"title": "Change", "body": "Run INSERT, UPDATE, or DELETE."},
                {"title": "Verify", "body": "Check affected rows and totals."},
                {"title": "Commit / rollback", "body": "Save or undo the work."},
            ], "caption": "The safest SQL professionals treat data change as a process, not as a single command."},
            {"type": "content", "title": "What DML is for", "bullets": [
                "`INSERT` adds new rows that did not previously exist.",
                "`UPDATE` changes existing rows that already exist.",
                "`DELETE` removes selected rows permanently from the table.",
                "Because these commands alter live data, they require more discipline than SELECT.",
            ], "note": [
                "This is the right place to remind learners that production systems often require approval and change controls.",
            ]},
            {"type": "code", "title": "INSERT patterns", "explanation": "Start with one row, then show how SQL can add many rows efficiently.", "code": [
                "INSERT INTO employees",
                "    (first_name, last_name, job_title, department_id, salary_zar, hire_date, email)",
                "VALUES",
                "    ('Simphiwe', 'Dube', 'Safety Officer', 4, 47500.00, '2024-05-01', 'si.dube@assmang.co.za');",
                "",
                "INSERT INTO employees (...) VALUES (...), (...), (...);",
            ], "takeaways": [
                "List columns explicitly so values cannot shift into the wrong positions.",
                "Multiple-row INSERT is more efficient than many separate statements.",
            ]},
            {"type": "code", "title": "UPDATE with a safety check", "explanation": "The safety SELECT is not optional in good teaching and good practice.", "code": [
                "SELECT employee_id, first_name, salary_zar",
                "FROM employees",
                "WHERE employee_id = 2;",
                "",
                "UPDATE employees",
                "SET salary_zar = 43000.00,",
                "    job_title = 'Senior HR Officer'",
                "WHERE employee_id = 2;",
            ], "takeaways": [
                "The WHERE clause defines the exact scope of the change.",
                "Verifying first prevents accidental bulk changes.",
            ]},
            {"type": "code", "title": "DELETE with verification", "explanation": "DELETE should be treated with the same discipline as UPDATE.", "code": [
                "SELECT *",
                "FROM training_register",
                "WHERE register_id = 14;",
                "",
                "DELETE FROM training_register",
                "WHERE register_id = 14;",
            ], "takeaways": [
                "The same WHERE clause appears in both the check and the delete.",
                "This keeps the operator mentally aligned with the intended scope.",
            ]},
            {"type": "comparison", "title": "DELETE vs TRUNCATE", "left_title": "DELETE", "left_points": [
                "Can remove selected rows with WHERE.",
                "Works naturally with transactions.",
                "More controlled for business operations.",
            ], "right_title": "TRUNCATE", "right_points": [
                "Removes all rows from the table.",
                "No WHERE clause.",
                "Treated as a much bigger structural action in practice.",
            ], "footer_text": "For beginner safety training, DELETE is a precise tool; TRUNCATE is a reset-style tool and should be used with caution."},
            {"type": "flow", "title": "Why transactions matter", "steps": [
                {"title": "START", "body": "Open a protected change window."},
                {"title": "Modify", "body": "Run one or several DML statements."},
                {"title": "Inspect", "body": "Check whether the results look correct."},
                {"title": "COMMIT", "body": "Save changes permanently."},
                {"title": "ROLLBACK", "body": "Undo the pending changes if needed."},
            ], "caption": "Transactions reduce fear and improve discipline because they encourage verification before final save."},
            {"type": "code", "title": "Transaction example", "explanation": "This pattern is ideal for salary adjustments and other high-impact changes.", "code": [
                "START TRANSACTION;",
                "",
                "UPDATE employees",
                "SET salary_zar = ROUND(salary_zar * 1.08, 2)",
                "WHERE department_id = 2;",
                "",
                "-- inspect the result here",
                "COMMIT;",
                "-- or ROLLBACK;",
            ], "takeaways": [
                "A transaction creates a deliberate review step.",
                "COMMIT finalises; ROLLBACK reverses pending work.",
            ]},
            {"type": "recap"},
        ],
    },
    {
        "day": "Day 2 • Topic 10",
        "title": "DDL — CREATE, ALTER, DROP and Best Practices",
        "slug": "topic-10-ddl-create-alter-drop-and-best-practices",
        "output": "Day-02/Topic-10-DDL-CREATE-ALTER-DROP-and-Best-Practices/topic-10-ddl-create-alter-drop-and-best-practices-presentation.pptx",
        "objectives": [
            "Explain the purpose of CREATE TABLE, ALTER TABLE, DROP TABLE, indexes, and views.",
            "Choose suitable beginner-level data types and constraints for relational tables.",
            "Read a CREATE TABLE statement as a schema definition, not just as syntax.",
            "Apply beginner SQL best practices for readability, safety, and maintainability.",
        ],
        "business_lens": [
            "New contractor tables.",
            "Adding columns safely.",
            "View-based reporting.",
            "Repeatable schema scripts.",
        ],
        "recap": [
            "DDL defines or changes structure, so its impact is broader than a single data change.",
            "Data types and constraints protect data quality before bad rows even enter the table.",
            "Readable SQL and safe change habits matter as much as syntax correctness.",
        ],
        "mistakes": [
            "Choosing vague or inappropriate data types.",
            "Dropping or altering objects without understanding downstream impact.",
            "Writing schema scripts that fail when rerun because IF EXISTS / IF NOT EXISTS was ignored.",
        ],
        "slides": [
            {"type": "objectives"},
            {"type": "content", "title": "What DDL changes", "bullets": [
                "DDL stands for Data Definition Language and focuses on structure rather than table contents.",
                "With DDL, SQL defines tables, columns, constraints, indexes, and views.",
                "A DDL script describes the shape of the system: what can be stored, what must be present, and how related tables connect.",
                "Because structure changes affect many users, DDL should be planned and reviewed carefully.",
            ], "note": [
                "Contrast this with DML: DML changes rows, DDL changes the container and rules.",
            ]},
            {"type": "matrix", "title": "Common data types and why they matter", "headers": ["Type", "Best beginner use", "Why it matters"], "rows": [
                ["INT", "Identifiers and whole-number counts", "Clear numeric storage for ids and counters."],
                ["VARCHAR(n)", "Names, titles, emails", "Stores text efficiently without fixed-length waste."],
                ["DECIMAL(p,s)", "Money and exact measured values", "Avoids floating-point surprises in financial reporting."],
                ["DATE", "Hire dates, contract dates", "Lets SQL reason about time correctly."],
                ["TINYINT(1)", "Boolean flags such as is_active", "Simple true/false style storage."],
            ], "note": "The key teaching point is not memorising types, but understanding what kind of value a column is meant to represent."},
            {"type": "code", "title": "Reading a CREATE TABLE statement", "explanation": "A CREATE TABLE statement is a business rules document written in SQL form.", "code": [
                "CREATE TABLE contractors (",
                "    contractor_id INT NOT NULL AUTO_INCREMENT,",
                "    company_name VARCHAR(150) NOT NULL,",
                "    contact_email VARCHAR(150) UNIQUE,",
                "    mine_id INT,",
                "    daily_rate_zar DECIMAL(10,2) NOT NULL DEFAULT 0.00,",
                "    PRIMARY KEY (contractor_id),",
                "    FOREIGN KEY (mine_id) REFERENCES mines(mine_id)",
                ");",
            ], "takeaways": [
                "The table defines names, types, defaults, and relationship rules in one place.",
                "Constraints protect the business from incomplete or inconsistent data.",
            ]},
            {"type": "nodes", "title": "Constraint thinking", "nodes": [
                {"id": "nn", "title": "NOT NULL", "x": Inches(0.8), "y": Inches(1.95), "w": Inches(2.3), "h": Inches(1.5), "fill": C_LIGHT_BLUE, "lines": ["Value must be supplied", "Good for essential fields"]},
                {"id": "pk", "title": "PRIMARY KEY", "x": Inches(3.45), "y": Inches(1.95), "w": Inches(2.3), "h": Inches(1.5), "fill": C_LIGHT_TEAL, "lines": ["Unique row identity", "Cannot be duplicated"]},
                {"id": "uq", "title": "UNIQUE", "x": Inches(6.1), "y": Inches(1.95), "w": Inches(2.3), "h": Inches(1.5), "fill": C_LIGHT_GOLD, "lines": ["No repeated values", "Useful for emails or codes"]},
                {"id": "fk", "title": "FOREIGN KEY", "x": Inches(8.75), "y": Inches(1.95), "w": Inches(2.8), "h": Inches(1.5), "fill": C_LIGHT_BLUE, "lines": ["Links to another table", "Protects referential integrity"]},
            ], "connections": [
                {"from": "pk", "to": "fk", "label": "reference path", "color": C_RED}
            ], "caption": "Constraints are business safeguards: they stop bad data from entering the system too easily."},
            {"type": "code", "title": "ALTER TABLE examples", "explanation": "ALTER TABLE evolves the schema when requirements change.", "code": [
                "ALTER TABLE contractors",
                "ADD COLUMN phone_number VARCHAR(20);",
                "",
                "ALTER TABLE contractors",
                "RENAME COLUMN is_active TO active_status;",
                "",
                "ALTER TABLE employees",
                "MODIFY COLUMN email VARCHAR(200) NOT NULL;",
            ], "takeaways": [
                "ALTER changes existing structure rather than creating a new table.",
                "Schema changes should be tested carefully because they affect later scripts and reports.",
            ]},
            {"type": "comparison", "title": "Views and indexes: different purposes", "left_title": "View", "left_points": [
                "A saved SELECT query that behaves like a virtual table.",
                "Improves reuse and simplifies reporting access.",
                "Good for safe presentation layers.",
            ], "right_title": "Index", "right_points": [
                "A structure that helps the database find rows faster.",
                "Improves performance for searched or joined columns.",
                "Does not change the visible result set.",
            ], "footer_text": "Views improve consumption; indexes improve retrieval speed."},
            {"type": "code", "title": "CREATE VIEW example", "explanation": "Views are powerful for beginner reporting because they hide complexity behind a reusable name.", "code": [
                "CREATE VIEW vw_employee_directory AS",
                "SELECT",
                "    e.employee_id,",
                "    CONCAT(e.first_name, ' ', e.last_name) AS full_name,",
                "    d.department_name",
                "FROM employees e",
                "INNER JOIN departments d ON e.department_id = d.department_id",
                "WHERE e.is_active = 1;",
            ], "takeaways": [
                "A view does not copy data; it stores the query definition.",
                "Business users can query the view without retyping the complex join every time.",
            ]},
            {"type": "content", "title": "SQL best practices to take forward", "bullets": [
                "Use consistent formatting so a reader can scan the query structure quickly.",
                "Choose explicit column lists over `SELECT *` in production-facing work.",
                "Name objects clearly and keep naming conventions predictable.",
                "Add IF EXISTS / IF NOT EXISTS where appropriate so scripts behave better when rerun.",
                "Use comments for intent, not to excuse unclear code.",
            ], "note": [
                "This is a strong closing slide before the overall course recap and final assessment.",
            ]},
            {"type": "recap"},
        ],
    },
]


def build_presentation(topic):
    prs = new_prs()
    title_slide(prs, topic["title"], topic["day"], topic["slug"])
    for slide_cfg in topic["slides"]:
        render_slide(prs, slide_cfg, topic)
    end_slide(prs, topic["title"])
    output = BASE / topic["output"]
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output, len(prs.slides)


def main():
    print("=" * 72)
    print("Generating detailed SQL100 presentations with diagrams...")
    print("=" * 72)
    for topic in TOPICS:
        output, count = build_presentation(topic)
        print(f"✅ {output.relative_to(BASE)}  ->  {count} slides")
    print("=" * 72)
    print(f"Completed {len(TOPICS)} presentation decks.")
    print("=" * 72)


if __name__ == "__main__":
    main()
